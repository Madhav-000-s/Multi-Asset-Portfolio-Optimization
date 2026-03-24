"""
Generate presentation.pptx — DSR Portfolio Optimizer Architecture Deep Dive

Run:  python create_presentation.py
Deps: pip install python-pptx matplotlib
"""

import io
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyArrowPatch, FancyBboxPatch
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Theme colours ─────────────────────────────────────────────────────────────
NAVY    = RGBColor(0x1f, 0x4e, 0x79)
BLUE    = RGBColor(0x2e, 0x86, 0xc1)
LBLUE   = RGBColor(0xae, 0xd6, 0xf1)
ORANGE  = RGBColor(0xe6, 0x7e, 0x22)   # R colour
YELLOW  = RGBColor(0xf4, 0xd0, 0x3f)   # Python colour
GREEN   = RGBColor(0x1e, 0x8b, 0x4c)
RED     = RGBColor(0xc0, 0x39, 0x2b)
WHITE   = RGBColor(0xff, 0xff, 0xff)
LIGHT   = RGBColor(0xf2, 0xf3, 0xf4)
DARK    = RGBColor(0x2c, 0x3e, 0x50)
GREY    = RGBColor(0x7f, 0x8c, 0x8d)

# hex versions for matplotlib
mNAVY   = "#1f4e79"
mBLUE   = "#2e86c1"
mLBLUE  = "#aed6f1"
mORANGE = "#e67e22"
mYELLOW = "#f4d03f"
mGREEN  = "#1e8b4c"
mLIGHT  = "#f2f3f4"
mDARK   = "#2c3e50"
mGREY   = "#7f8c8d"
mRED    = "#c0392b"

SLD_W = 10   # slide width inches
SLD_H = 7.5  # slide height inches


# ── Helpers ───────────────────────────────────────────────────────────────────

def _rgb(r):
    return RGBColor(r[0], r[1], r[2])

def _set_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def _add_textbox(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=WHITE,
                 align=PP_ALIGN.LEFT, word_wrap=True):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txb.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb

def _add_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=Pt(1)):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def _add_rect_text(slide, text, left, top, width, height,
                   fill_color, text_color=WHITE, font_size=14,
                   bold=False, line_color=None, align=PP_ALIGN.CENTER):
    shape = _add_rect(slide, left, top, width, height, fill_color, line_color)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = text_color
    return shape

def _title_slide_header(slide, title, subtitle=None):
    _add_rect(slide, 0, 0, SLD_W, SLD_H, NAVY)
    _add_rect(slide, 0, 0, SLD_W, 0.08, ORANGE)          # top accent bar
    _add_rect(slide, 0, SLD_H - 0.08, SLD_W, 0.08, ORANGE)  # bottom accent bar
    _add_textbox(slide, title, 0.6, 1.5, 8.8, 2.5,
                 font_size=36, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER)
    if subtitle:
        _add_textbox(slide, subtitle, 0.6, 4.2, 8.8, 1.2,
                     font_size=20, color=LBLUE, align=PP_ALIGN.CENTER)

def _section_header(slide, title, subtitle=None, accent=NAVY):
    _add_rect(slide, 0, 0, SLD_W, 1.35, accent)
    _add_textbox(slide, title, 0.4, 0.18, 9.2, 0.8,
                 font_size=26, bold=True, color=WHITE)
    if subtitle:
        _add_textbox(slide, subtitle, 0.4, 0.9, 9.2, 0.45,
                     font_size=14, color=LBLUE)

def _png_to_slide(slide, fig, left, top, width, height):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight",
                facecolor=fig.get_facecolor(), dpi=150)
    buf.seek(0)
    slide.shapes.add_picture(buf, Inches(left), Inches(top),
                             Inches(width), Inches(height))
    plt.close(fig)

def _bullet_list(slide, items, left, top, width, height,
                 font_size=16, color=DARK, spacing=0.38):
    for i, item in enumerate(items):
        indent = item.startswith("  ")
        text = item.lstrip()
        prefix = "    •  " if indent else "•  "
        _add_textbox(slide, prefix + text,
                     left + (0.3 if indent else 0),
                     top + i * spacing,
                     width - (0.3 if indent else 0),
                     spacing + 0.05,
                     font_size=font_size - (1 if indent else 0),
                     color=color)


# ══════════════════════════════════════════════════════════════════════════════
# DIAGRAM GENERATORS
# ══════════════════════════════════════════════════════════════════════════════

def diag_system_architecture():
    fig, ax = plt.subplots(figsize=(11, 7))
    fig.patch.set_facecolor("#f8f9fa")
    ax.set_xlim(0, 11); ax.set_ylim(0, 7); ax.axis("off")

    layers = [
        ("Layer 1", "Data Ingestion",       "yfinance · Alpha Vantage · NewsAPI",   mBLUE,   6.2),
        ("Layer 2", "Feature Engineering",  "Log Ret · Vol · RSI · MACD · FinBERT", "#5d6d7e", 4.8),
        ("Layer 3", "LSTM + DSR Loss",      "PyTorch · 2-Layer LSTM · Softmax",     mNAVY,   3.4),
        ("Layer 4", "R Risk Engine",        "PortfolioAnalytics · PerformanceAnalytics", mORANGE, 2.0),
        ("Layer 5", "Dashboard & Report",   "Streamlit · R Shiny · R Markdown",     mGREEN,  0.6),
    ]

    for (lnum, lname, ldetail, lcolor, ly) in layers:
        # main box
        box = FancyBboxPatch((0.3, ly), 7.5, 1.1,
                             boxstyle="round,pad=0.07", linewidth=2,
                             edgecolor=lcolor, facecolor=lcolor + "22")
        ax.add_patch(box)
        ax.text(0.7, ly + 0.75, lnum, fontsize=9, color=lcolor,
                fontweight="bold", va="center")
        ax.text(0.7, ly + 0.45, lname, fontsize=14, color=mDARK,
                fontweight="bold", va="center")
        ax.text(0.7, ly + 0.18, ldetail, fontsize=9.5, color=mGREY, va="center")

        # artifact box on the right
        artifacts = {
            "Layer 1": "prices.parquet\nheadlines.csv",
            "Layer 2": "features (30–35 dims)\nsentiment_scores.parquet",
            "Layer 3": "best_model.pt\nweights w_t ∈ ℝ⁵",
            "Layer 4": "metrics.csv\nweights_history.csv",
            "Layer 5": "Interactive Dashboard\nHTML/PDF Report",
        }
        art = FancyBboxPatch((8.1, ly + 0.1), 2.6, 0.88,
                             boxstyle="round,pad=0.05", linewidth=1.2,
                             edgecolor=lcolor, facecolor="#ffffff")
        ax.add_patch(art)
        ax.text(9.4, ly + 0.54, artifacts[lnum],
                fontsize=8.5, color=mDARK, ha="center", va="center",
                linespacing=1.5)
        ax.annotate("", xy=(8.1, ly + 0.54), xytext=(7.8, ly + 0.54),
                    arrowprops=dict(arrowstyle="->", color=lcolor, lw=1.5))

        # downward arrow between layers
        if ly > 0.6:
            ax.annotate("", xy=(4.05, ly), xytext=(4.05, ly + 1.1),
                        arrowprops=dict(arrowstyle="->", color=mGREY, lw=1.8))

    # reticulate label — positioned in the gap between Layer 3 (y=3.4) and Layer 4 (top y=3.1)
    ax.text(4.4, 3.2, "reticulate", fontsize=8.5, color=mORANGE,
            style="italic", ha="center",
            bbox=dict(facecolor="white", edgecolor=mORANGE, pad=2, lw=1))

    ax.set_title("System Architecture — 5-Layer Pipeline",
                 fontsize=15, fontweight="bold", color=mDARK, pad=10)
    return fig


def diag_r_components():
    fig, ax = plt.subplots(figsize=(11, 6))
    fig.patch.set_facecolor("#fff8f2")
    ax.set_xlim(0, 11); ax.set_ylim(0, 6); ax.axis("off")

    # Central orchestrator
    orch = FancyBboxPatch((3.8, 4.0), 3.4, 1.0,
                          boxstyle="round,pad=0.1", lw=2.5,
                          edgecolor=mORANGE, facecolor="#fdebd0")
    ax.add_patch(orch)
    ax.text(5.5, 4.75, "backtest.R", fontsize=13, fontweight="bold",
            color=mORANGE, ha="center")
    ax.text(5.5, 4.25, "Walk-forward orchestrator", fontsize=9,
            color=mGREY, ha="center")

    # Sub-modules
    modules = [
        (0.3, 2.2, "constraints.R", "Min/max bounds\nCVaR optimisation", "#d35400"),
        (2.7, 2.2, "transaction_costs.R", "Turnover calc\n10 bps deduction", "#e67e22"),
        (5.1, 2.2, "metrics.R", "Sharpe · Sortino\nVaR · CVaR · Calmar", "#2e86c1"),
        (7.5, 2.2, "efficient_frontier.R", "Mean-variance\nfrontier plot", "#1e8b4c"),
    ]
    for (mx, my, mname, mdesc, mcol) in modules:
        b = FancyBboxPatch((mx, my), 2.2, 1.3,
                           boxstyle="round,pad=0.07", lw=1.8,
                           edgecolor=mcol, facecolor=mcol + "22")
        ax.add_patch(b)
        ax.text(mx + 1.1, my + 0.95, mname, fontsize=9.5, fontweight="bold",
                color=mcol, ha="center")
        ax.text(mx + 1.1, my + 0.45, mdesc, fontsize=8.5, color=mDARK,
                ha="center", va="center", linespacing=1.4)
        # Arrow from backtest.R down to module
        ax.annotate("", xy=(mx + 1.1, my + 1.3), xytext=(5.5, 4.0),
                    arrowprops=dict(arrowstyle="->", color=mORANGE,
                                   lw=1.3, connectionstyle="arc3,rad=0.0"))

    # Python on the right
    py_box = FancyBboxPatch((8.8, 4.0), 2.0, 1.0,
                            boxstyle="round,pad=0.08", lw=2,
                            edgecolor="#3498db", facecolor="#ebf5fb")
    ax.add_patch(py_box)
    ax.text(9.8, 4.74, "predict.py", fontsize=11, fontweight="bold",
            color="#2980b9", ha="center")
    ax.text(9.8, 4.25, "(Python / PyTorch)", fontsize=8, color=mGREY, ha="center")
    ax.annotate("", xy=(8.8, 4.5), xytext=(7.2, 4.5),
                arrowprops=dict(arrowstyle="<->", color="#3498db", lw=1.8))
    ax.text(8.0, 4.65, "reticulate", fontsize=8, color="#3498db",
            ha="center", style="italic")

    # Outputs at bottom
    outputs = [
        (1.5, 0.2, "weights_history.csv", mBLUE),
        (4.0, 0.2, "metrics.csv", mBLUE),
        (6.5, 0.2, "backtest_results.rds", mBLUE),
        (8.8, 0.2, "*.png plots", mBLUE),
    ]
    for (ox, oy, oname, ocol) in outputs:
        ob = FancyBboxPatch((ox, oy), 2.2, 0.7,
                            boxstyle="round,pad=0.05", lw=1,
                            edgecolor=ocol, facecolor="#eaf2ff")
        ax.add_patch(ob)
        ax.text(ox + 1.1, oy + 0.35, oname, fontsize=8.5, color=mDARK,
                ha="center", va="center")

    ax.annotate("", xy=(5.5, 0.9), xytext=(5.5, 2.2),
                arrowprops=dict(arrowstyle="->", color=mGREY, lw=1.5))
    ax.text(5.9, 1.5, "Saved\noutputs", fontsize=8, color=mGREY)

    ax.set_title("R Layer 4 — Component Diagram",
                 fontsize=14, fontweight="bold", color=mORANGE, pad=10)
    return fig


def diag_sequence_reticulate():
    fig, ax = plt.subplots(figsize=(11, 6.5))
    fig.patch.set_facecolor("#f8f9fa")
    ax.set_xlim(0, 11); ax.set_ylim(0, 6.5); ax.axis("off")

    # Lifelines
    lifelines = [
        (1.5,  "backtest.R\n(R)",       mORANGE, "#fdebd0"),
        (5.5,  "reticulate\n(bridge)",  mGREY,   "#f2f3f4"),
        (9.5,  "predict.py\n(Python)",  "#2980b9","#ebf5fb"),
    ]
    for (lx, lname, lcol, lfill) in lifelines:
        box = FancyBboxPatch((lx - 1.0, 5.6), 2.0, 0.75,
                             boxstyle="round,pad=0.07", lw=2,
                             edgecolor=lcol, facecolor=lfill)
        ax.add_patch(box)
        ax.text(lx, 5.97, lname, fontsize=10, fontweight="bold",
                color=lcol, ha="center", va="center", linespacing=1.3)
        ax.plot([lx, lx], [0.2, 5.6], color=lcol, lw=1.2,
                linestyle="--", alpha=0.5)

    # Messages
    messages = [
        (5.2, "source_python('predict.py')",    1.5, 5.5, "->",  mORANGE),
        (4.6, "predict_at_date('2024-01-08')",  1.5, 5.5, "->",  mORANGE),
        (4.3, "forward call",                   5.5, 9.5, "->",  mGREY),
        (3.8, "  load_model() — if not cached", 9.5, 9.5, "",    "#2980b9"),
        (3.3, "  compute_features(prices)",     9.5, 9.5, "",    "#2980b9"),
        (2.8, "  lstm.forward(X)",              9.5, 9.5, "",    "#2980b9"),
        (2.3, "return np.array([w1…w5])",       9.5, 5.5, "<-",  "#2980b9"),
        (1.8, "weights as R vector",            5.5, 1.5, "<-",  mGREY),
        (1.3, "apply_constraints(weights)",     1.5, 1.5, "",    mORANGE),
        (0.8, "w · r_{t+1}  →  portfolio return",1.5,1.5,"",   mORANGE),
    ]
    for (my, msg, x1, x2, arrow, col) in messages:
        if arrow == "->":
            ax.annotate("", xy=(x2 - 0.05, my), xytext=(x1 + 0.05, my),
                        arrowprops=dict(arrowstyle="->", color=col, lw=1.5))
            ax.text((x1 + x2) / 2, my + 0.13, msg,
                    fontsize=8.5, color=col, ha="center")
        elif arrow == "<-":
            ax.annotate("", xy=(x2 + 0.05, my), xytext=(x1 - 0.05, my),
                        arrowprops=dict(arrowstyle="->", color=col, lw=1.5))
            ax.text((x1 + x2) / 2, my + 0.13, msg,
                    fontsize=8.5, color=col, ha="center")
        else:
            ax.text(x1 + 0.15, my, msg, fontsize=8.5, color=col, va="center")

    # Loop box
    loop = plt.Rectangle((0.4, 0.4), 6.8, 4.2,
                          fill=False, edgecolor=mORANGE,
                          linestyle=":", lw=1.5)
    ax.add_patch(loop)
    ax.text(0.55, 4.7, "loop  [each weekly rebalance date]",
            fontsize=8, color=mORANGE, style="italic")

    ax.set_title("R ↔ Python Sequence Diagram (via reticulate)",
                 fontsize=14, fontweight="bold", color=mDARK, pad=10)
    return fig


def diag_lstm_architecture():
    fig, ax = plt.subplots(figsize=(11, 4.5))
    fig.patch.set_facecolor("#f0f4f8")
    ax.set_xlim(0, 11); ax.set_ylim(0, 4.5); ax.axis("off")

    # 5 blocks, width=1.7 each, gap=0.52, left margin=0.2
    blocks = [
        (0.2,  "Input\n(batch, 60, 30)\n6 features × 5 assets", mBLUE,   1.7),
        (2.42, "LSTM Layer 1\nhidden = 128\ndropout = 0.3",      mNAVY,   1.7),
        (4.64, "LSTM Layer 2\nhidden = 64\ndropout = 0.3",       mNAVY,   1.7),
        (6.86, "Linear\n64 → 5",                                 "#5d6d7e",1.7),
        (9.08, "Softmax\nweights ∈ ℝ⁵\nsum = 1",                mGREEN,  1.72),
    ]
    for (bx, btext, bcol, bw) in blocks:
        b = FancyBboxPatch((bx, 1.1), bw, 2.3,
                           boxstyle="round,pad=0.1", lw=2,
                           edgecolor=bcol, facecolor=bcol + "22")
        ax.add_patch(b)
        ax.text(bx + bw/2, 2.25, btext, fontsize=9, color=mDARK,
                ha="center", va="center", linespacing=1.5, fontweight="bold")
        if bx < 9.0:
            ax.annotate("", xy=(bx + bw + 0.42, 2.25),
                        xytext=(bx + bw + 0.08, 2.25),
                        arrowprops=dict(arrowstyle="->", color=mGREY, lw=2))

    # DSR Loss annotation
    ax.text(5.5, 0.55, "← Trained with Differential Sharpe Ratio (DSR) Loss — not cross-entropy! →",
            fontsize=9, color=mORANGE, ha="center", style="italic")
    ax.text(5.5, 0.2, "η = 0.01  |  Adam optimiser  |  Early stopping on val DSR",
            fontsize=8.5, color=mGREY, ha="center")

    ax.set_title("LSTM Portfolio Model Architecture",
                 fontsize=14, fontweight="bold", color=mDARK, pad=8)
    return fig


def diag_cvar_flow():
    fig, ax = plt.subplots(figsize=(10, 5))
    fig.patch.set_facecolor("#fff8f2")
    ax.set_xlim(0, 10); ax.set_ylim(0, 5); ax.axis("off")

    steps = [
        (0.3, 3.5, "LSTM Output\nraw softmax\nweights w_t",     mBLUE),
        (2.8, 3.5, "Box Constraints\nw_min = 2%\nw_max = 30%",  mORANGE),
        (5.3, 3.5, "CVaR Check\n95% confidence\nlevel",         mRED),
        (7.8, 3.5, "Final Portfolio\nweights\n∑w = 1",          mGREEN),
    ]
    for (sx, sy, stext, scol) in steps:
        b = FancyBboxPatch((sx, sy), 1.9, 1.3,
                           boxstyle="round,pad=0.1", lw=2,
                           edgecolor=scol, facecolor=scol + "22")
        ax.add_patch(b)
        ax.text(sx + 0.95, sy + 0.65, stext,
                fontsize=9.5, color=mDARK, ha="center", va="center",
                linespacing=1.5, fontweight="bold")
        if sx < 7.8:
            ax.annotate("", xy=(sx + 2.05, sy + 0.65),
                        xytext=(sx + 1.9, sy + 0.65),
                        arrowprops=dict(arrowstyle="->", color=mGREY, lw=2))

    # R function labels
    labels = [
        (1.25, 3.3, "predict_at_date()", mBLUE),
        (3.75, 3.3, "constrain_weights()", mORANGE),
        (6.25, 3.3, "compute_portfolio_cvar()", mRED),
    ]
    for (lx, ly, ltxt, lcol) in labels:
        ax.text(lx, ly, ltxt, fontsize=7.5, color=lcol,
                ha="center", style="italic")

    # Rejection loop
    ax.annotate("", xy=(6.25, 4.8), xytext=(6.25, 3.5 + 1.3),
                arrowprops=dict(arrowstyle="->", color=mRED, lw=1.3))
    ax.annotate("", xy=(3.75, 4.8), xytext=(6.25, 4.8),
                arrowprops=dict(arrowstyle="->", color=mRED, lw=1.3,
                                connectionstyle="arc3,rad=0"))
    ax.annotate("", xy=(3.75, 4.8), xytext=(3.75, 3.5 + 1.3),
                arrowprops=dict(arrowstyle="->", color=mRED, lw=1.3))
    ax.text(5.0, 4.95, "CVaR violated → re-optimise",
            fontsize=8, color=mRED, ha="center", style="italic")

    # Transaction cost note
    ax.text(5.0, 2.5, "Transaction cost: 10 bps × |Δw|  applied at every rebalance  (backtest.R · transaction_costs.R)",
            fontsize=9, color=mGREY, ha="center")

    ax.set_title("R: Weight Constraint & CVaR Optimisation Flow",
                 fontsize=13, fontweight="bold", color=mORANGE, pad=10)
    return fig


def diag_data_flow():
    fig, ax = plt.subplots(figsize=(11, 5.5))
    fig.patch.set_facecolor("#f8f9fa")
    ax.set_xlim(0, 11); ax.set_ylim(0, 5.5); ax.axis("off")

    nodes = [
        # (x, y, label, sublabel, colour)
        (0.1, 2.5, "config.yaml",         "assets, dates\nhyperparams",  "#5d6d7e"),
        (1.9, 4.2, "data_loader.py",      "yfinance\nOHLCV",             mBLUE),
        (1.9, 0.8, "scraper.py",          "Alpha Vantage\nheadlines",    mBLUE),
        (4.0, 4.2, "features.py",         "30 tech\nfeatures",           mNAVY),
        (4.0, 0.8, "sentiment.py",        "FinBERT\nscores",             mNAVY),
        (6.1, 2.5, "train.py",            "LSTM +\nDSR loss",            mNAVY),
        (8.1, 2.5, "backtest.R",          "walk-forward\nR engine",      mORANGE),
        (10.0, 4.2, "Dashboard",          "Streamlit\nShiny",            mGREEN),
        (10.0, 0.8, "Report",             "R Markdown\nHTML/PDF",        mGREEN),
    ]
    for (nx, ny, nl, nsl, nc) in nodes:
        b = FancyBboxPatch((nx, ny), 1.6, 1.1,
                           boxstyle="round,pad=0.07", lw=1.8,
                           edgecolor=nc, facecolor=nc + "22")
        ax.add_patch(b)
        ax.text(nx + 0.8, ny + 0.75, nl, fontsize=8.5, fontweight="bold",
                color=mDARK, ha="center")
        ax.text(nx + 0.8, ny + 0.28, nsl, fontsize=7.5, color=mGREY,
                ha="center", linespacing=1.3)

    # Edges
    edges = [
        (0.1+1.6, 3.05, 1.9, 4.7),       # config → data_loader
        (0.1+1.6, 3.05, 1.9, 1.35),      # config → scraper
        (1.9+1.6, 4.7, 4.0, 4.7),        # data_loader → features
        (1.9+1.6, 1.35, 4.0, 1.35),      # scraper → sentiment
        (4.0+1.6, 4.7, 6.1+0.8, 3.6),    # features → train
        (4.0+1.6, 1.35, 6.1+0.8, 2.5),   # sentiment → train
        (6.1+1.6, 3.05, 8.1, 3.05),      # train → backtest
        (8.1+1.6, 4.1, 10.0, 4.7),       # backtest → dashboard
        (8.1+1.6, 2.2, 10.0, 1.35),      # backtest → report
    ]
    for (x1, y1, x2, y2) in edges:
        ax.annotate("", xy=(x2, y2), xytext=(x1, y1),
                    arrowprops=dict(arrowstyle="->", color=mGREY,
                                   lw=1.4, connectionstyle="arc3,rad=0.15"))

    ax.text(5.5, 0.15, "Python  ←————————————→  R",
            fontsize=9, color=mGREY, ha="center", style="italic")
    ax.axvline(7.5, color=mORANGE, lw=1.2, linestyle="--", alpha=0.4)
    ax.text(7.6, 5.2, "R boundary", fontsize=8, color=mORANGE, style="italic")

    ax.set_title("End-to-End Data Flow",
                 fontsize=14, fontweight="bold", color=mDARK, pad=8)
    return fig


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ══════════════════════════════════════════════════════════════════════════════

def slide_title(prs):
    layout = prs.slide_layouts[6]   # blank
    sld = prs.slides.add_slide(layout)
    _set_bg(sld, NAVY)
    _add_rect(sld, 0, 0, SLD_W, 0.12, ORANGE)
    _add_rect(sld, 0, SLD_H - 0.12, SLD_W, 0.12, ORANGE)
    # Decorative strip
    _add_rect(sld, 0, 2.8, 0.25, 2.0, ORANGE)
    _add_textbox(sld,
        "Deep Portfolio Optimization\nUsing DSR + Sentiment-Augmented LSTM",
        0.6, 0.9, 8.8, 2.4,
        font_size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _add_textbox(sld, "Architecture Deep Dive", 0.6, 3.4, 8.8, 0.7,
                 font_size=22, color=LBLUE, align=PP_ALIGN.CENTER)
    _add_textbox(sld, "With a special focus on R's role in the pipeline",
                 0.6, 4.1, 8.8, 0.6,
                 font_size=14, color=GREY, align=PP_ALIGN.CENTER)
    _add_textbox(sld, "Python  ·  R  ·  PyTorch  ·  reticulate",
                 0.6, 5.6, 8.8, 0.6,
                 font_size=12, color=LBLUE, align=PP_ALIGN.CENTER)


def slide_what_are_we_building(prs):
    layout = prs.slide_layouts[6]
    sld = prs.slides.add_slide(layout)
    _set_bg(sld, LIGHT)
    _section_header(sld, "What Are We Building?",
                    "A casual overview — no jargon required")

    _add_textbox(sld, "The big idea:", 0.5, 1.5, 4.0, 0.5,
                 font_size=14, bold=True, color=NAVY)
    _add_textbox(sld,
        '"Use deep learning to figure out how to split money across 5 big tech stocks, '
        'then let R double-check that we\'re not taking crazy risks."',
        0.5, 1.9, 5.5, 1.5,
        font_size=13, color=DARK)

    bullets = [
        "We have 5 assets: AAPL, MSFT, GOOGL, AMZN, META",
        "An LSTM reads 60 days of price history each week",
        "It spits out 5 weights  (e.g. 30% Apple, 20% Google…)",
        "R then checks: are those weights safe? Too concentrated?",
        "R runs the backtest and measures how good we'd have done",
        "A Streamlit dashboard shows it all live",
    ]
    _bullet_list(sld, bullets, 0.5, 3.55, 6.0, 0.45, font_size=14, color=DARK)

    # right callout
    _add_rect_text(sld, "Training data\n2018 – 2024", 7.0, 1.6, 2.6, 1.0,
                   fill_color=NAVY, font_size=13)
    _add_rect_text(sld, "Test period\nDec 2023 – Dec 2024", 7.0, 2.8, 2.6, 1.0,
                   fill_color=BLUE, font_size=13)
    _add_rect_text(sld, "Rebalance\nEvery week · 10 bps cost", 7.0, 4.0, 2.6, 1.0,
                   fill_color=_rgb((0x2e,0x86,0x4c)), font_size=13)


def slide_architecture(prs):
    layout = prs.slide_layouts[6]
    sld = prs.slides.add_slide(layout)
    _set_bg(sld, LIGHT)
    _section_header(sld, "System Architecture — 5 Layers",
                    "Each layer feeds the next; Python hands off to R at Layer 4")
    fig = diag_system_architecture()
    _png_to_slide(sld, fig, 0.3, 1.3, 9.4, 6.0)


def slide_data_flow(prs):
    layout = prs.slide_layouts[6]
    sld = prs.slides.add_slide(layout)
    _set_bg(sld, LIGHT)
    _section_header(sld, "Data Flow — End to End",
                    "From raw Yahoo Finance data to an interactive dashboard")
    fig = diag_data_flow()
    _png_to_slide(sld, fig, 0.2, 1.3, 9.6, 5.9)


def slide_layer1(prs):
    layout = prs.slide_layouts[6]
    sld = prs.slides.add_slide(layout)
    _set_bg(sld, LIGHT)
    _section_header(sld, "Layer 1 — Data Ingestion",
                    "Pulling price history and news headlines from the web",
                    accent=BLUE)

    _add_rect_text(sld, "Prices (yfinance)", 0.4, 1.5, 2.8, 0.7,
                   fill_color=BLUE, font_size=14)
    bullets_p = [
        "Daily OHLCV for 5 tickers",
        "2018-01-01 → 2024-12-31",
        "Saved as prices.parquet",
        "Delta-updated for live mode",
    ]
    _bullet_list(sld, bullets_p, 0.4, 2.3, 3.2, 0.42, font_size=13, color=DARK)

    _add_rect_text(sld, "News Headlines (Alpha Vantage)", 4.0, 1.5, 3.2, 0.7,
                   fill_color=ORANGE, font_size=13)
    bullets_n = [
        "Historical articles per ticker",
        "Free API — 25 req/day",
        "Saved as headlines.csv",
        "Also yfinance recent news",
    ]
    _bullet_list(sld, bullets_n, 4.0, 2.3, 3.2, 0.42, font_size=13, color=DARK)

    _add_rect_text(sld, "config.yaml", 7.5, 1.5, 2.1, 0.7,
                   fill_color=NAVY, font_size=14)
    bullets_c = [
        "Asset universe",
        "Date range",
        "Hyper-parameters",
        "Constraint limits",
    ]
    _bullet_list(sld, bullets_c, 7.5, 2.3, 2.1, 0.42, font_size=13, color=DARK)

    _add_textbox(sld,
        "Why parquet?  Columnar format → 10× faster than CSV for financial time-series. "
        "Delta updates append only new rows → no re-downloading 6 years of data every refresh.",
        0.4, 5.6, 9.2, 0.9,
        font_size=12, color=GREY)


def slide_layer2(prs):
    layout = prs.slide_layouts[6]
    sld = prs.slides.add_slide(layout)
    _set_bg(sld, LIGHT)
    _section_header(sld, "Layer 2 — Feature Engineering",
                    "Turning raw prices into a rich 3D tensor the LSTM can read",
                    accent=_rgb((0x5d,0x6d,0x7e)))

    # Feature table
    features = [
        ("Log Return",     "ln(P_t / P_{t-1})",             "Daily price momentum"),
        ("Rolling Vol",    "20d std × √252",                 "Annualised volatility"),
        ("RSI-14",         "Relative Strength Index",        "Overbought / oversold"),
        ("MACD",           "EMA(12) − EMA(26) signal",       "Trend direction"),
        ("Bollinger %B",   "(Price − lower) / (upper−lower)", "Band position"),
        ("SMA Ratio",      "SMA(50) / SMA(200)",             "Golden/death cross"),
        ("Sentiment *",    "P(pos) − P(neg)  via FinBERT",  "News alpha signal"),
    ]
    headers = ["Feature", "Formula", "What it captures"]
    col_x   = [0.3, 3.0, 6.2]
    col_w   = [2.6, 3.1, 3.5]

    for ci, h in enumerate(headers):
        _add_rect_text(sld, h, col_x[ci], 1.45, col_w[ci], 0.42,
                       fill_color=NAVY, font_size=12, bold=True)

    for ri, (f, formula, meaning) in enumerate(features):
        y = 1.9 + ri * 0.52
        bg = LIGHT if ri % 2 == 0 else WHITE
        for ci, val in enumerate([f, formula, meaning]):
            _add_rect_text(sld, val, col_x[ci], y, col_w[ci], 0.48,
                           fill_color=bg, text_color=DARK,
                           font_size=11, line_color=GREY)

    _add_textbox(sld,
        "* Sentiment is the 7th feature — optional, enabled via config.yaml. "
        "Output tensor: (batch, 60 days, 30 or 35 features)  ·  Z-score normalised with 252-day rolling window.",
        0.3, 5.65, 9.4, 0.8,
        font_size=11, color=GREY)


def slide_layer3(prs):
    layout = prs.slide_layouts[6]
    sld = prs.slides.add_slide(layout)
    _set_bg(sld, LIGHT)
    _section_header(sld, "Layer 3 — LSTM + Differential Sharpe Ratio",
                    "The core ML engine — trained to maximise risk-adjusted returns directly",
                    accent=NAVY)
    fig = diag_lstm_architecture()
    _png_to_slide(sld, fig, 0.3, 1.35, 9.4, 3.7)

    _add_textbox(sld, "What makes this different from regular deep learning?",
                 0.4, 5.1, 9.2, 0.45, font_size=13, bold=True, color=NAVY)
    _add_textbox(sld,
        "Normal neural nets minimise prediction error (MSE).  "
        "This model minimises −DSR: the negative Differential Sharpe Ratio.  "
        "That means it directly learns to maximise risk-adjusted profit — no intermediate step.",
        0.4, 5.55, 9.2, 0.9,
        font_size=12, color=DARK)


def slide_r_overview(prs):
    layout = prs.slide_layouts[6]
    sld = prs.slides.add_slide(layout)
    _set_bg(sld, _rgb((0xff, 0xf8, 0xf2)))
    _section_header(sld, "Layer 4 — R Risk Engine: Overview",
                    "R doesn't just make charts — it does the heavy quantitative lifting",
                    accent=ORANGE)
    fig = diag_r_components()
    _png_to_slide(sld, fig, 0.2, 1.35, 9.6, 5.8)


def slide_reticulate(prs):
    layout = prs.slide_layouts[6]
    sld = prs.slides.add_slide(layout)
    _set_bg(sld, _rgb((0xff, 0xf8, 0xf2)))
    _section_header(sld, "R ↔ Python Bridge: reticulate",
                    "R calls the LSTM model directly — no files exchanged, no REST API needed",
                    accent=ORANGE)
    fig = diag_sequence_reticulate()
    _png_to_slide(sld, fig, 0.2, 1.35, 9.6, 5.8)


def slide_backtest(prs):
    layout = prs.slide_layouts[6]
    sld = prs.slides.add_slide(layout)
    _set_bg(sld, _rgb((0xff, 0xf8, 0xf2)))
    _section_header(sld, "R's Role: Walk-Forward Backtesting",
                    "backtest.R — 600 lines, the backbone of the whole experiment",
                    accent=ORANGE)

    steps = [
        ("1", "Load Data",      "Read prices.parquet\nvia arrow package",       mBLUE),
        ("2", "Get Dates",      "Weekly rebalance\ndates (xts endpoints)",      mNAVY),
        ("3", "Predict",        "Call Python:\npredict_at_date()",               "#2980b9"),
        ("4", "Constrain",      "Apply min/max\n& CVaR bounds",                 mORANGE),
        ("5", "Return",         "w_t · r_{t+1}\nminus 10 bps cost",            mGREEN),
        ("6", "Save",           "weights_history.csv\nmetrics.csv  .rds",      mGREY),
    ]
    for i, (num, title, desc, col) in enumerate(steps):
        x = 0.3 + i * 1.6
        _add_rect_text(sld, f"Step {num}\n{title}", x, 1.55, 1.4, 0.9,
                       fill_color=_rgb(
                           (int(col[1:3],16), int(col[3:5],16), int(col[5:7],16))
                       ) if col.startswith("#") else NAVY,
                       font_size=11, bold=True)
        _add_rect_text(sld, desc, x, 2.55, 1.4, 1.1,
                       fill_color=LIGHT, text_color=DARK,
                       font_size=10, line_color=GREY)
        if i < 5:
            _add_textbox(sld, "→", x + 1.42, 1.9, 0.2, 0.5,
                         font_size=18, color=GREY)

    _add_textbox(sld,
        "xts time series in R keeps everything perfectly date-aligned. "
        "PerformanceAnalytics::Return.portfolio() handles weight × return multiplication correctly across splits.",
        0.3, 4.0, 9.4, 0.9, font_size=12, color=DARK)

    bullets = [
        "~150 rebalance points over 2023–2024 test window",
        "Each point: Python inference (~0.3s) + R constraint check + return calculation",
        "Total backtest runtime: ~2 minutes on CPU",
    ]
    _bullet_list(sld, bullets, 0.3, 4.95, 9.4, 0.42, font_size=13, color=DARK)


def slide_cvar(prs):
    layout = prs.slide_layouts[6]
    sld = prs.slides.add_slide(layout)
    _set_bg(sld, _rgb((0xff, 0xf8, 0xf2)))
    _section_header(sld, "R's Role: PortfolioAnalytics & CVaR",
                    "Turning raw LSTM weights into regulated, risk-controlled positions",
                    accent=ORANGE)
    fig = diag_cvar_flow()
    _png_to_slide(sld, fig, 0.3, 1.4, 9.4, 4.5)

    _add_textbox(sld,
        "Why PortfolioAnalytics?  It provides an ROI solver specifically designed for portfolio "
        "constraints — min/max weights, CVaR limits, long-only, sum-to-one — all in one call. "
        "scipy.optimize can do it, but PortfolioAnalytics handles the edge cases (infeasible regions, etc.) out of the box.",
        0.3, 6.0, 9.4, 1.1, font_size=11, color=GREY)


def slide_performance_analytics(prs):
    layout = prs.slide_layouts[6]
    sld = prs.slides.add_slide(layout)
    _set_bg(sld, _rgb((0xff, 0xf8, 0xf2)))
    _section_header(sld, "R's Role: PerformanceAnalytics",
                    "8 industry-standard risk metrics computed in <10 lines of R",
                    accent=ORANGE)

    metrics = [
        ("Annualised Return",    "Return.annualized(r, scale=252)",         "How much did we make?"),
        ("Sharpe Ratio",         "SharpeRatio.annualized(r, Rf=0)",         "Return per unit of risk"),
        ("Sortino Ratio",        "SortinoRatio(r)",                         "Return per unit of downside risk"),
        ("Max Drawdown",         "maxDrawdown(r)",                           "Worst peak-to-trough loss"),
        ("Calmar Ratio",         "Return / maxDrawdown(r)",                  "Return vs max pain"),
        ("VaR (95%)",            "VaR(r, p=0.95, method='historical')",     "Daily loss at 95% confidence"),
        ("CVaR / ES (95%)",      "ES(r, p=0.95, method='historical')",      "Expected loss beyond VaR"),
        ("Annualised Volatility","StdDev.annualized(r, scale=252)",         "How bumpy was the ride?"),
    ]
    col_x = [0.3, 3.0, 6.4]; col_w = [2.6, 3.3, 3.3]
    headers = ["Metric", "R Function (PerformanceAnalytics)", "Plain English"]
    for ci, h in enumerate(headers):
        _add_rect_text(sld, h, col_x[ci], 1.45, col_w[ci], 0.38,
                       fill_color=ORANGE, font_size=11, bold=True)
    for ri, (m, fn, pe) in enumerate(metrics):
        y = 1.86 + ri * 0.48
        bg = LIGHT if ri % 2 == 0 else WHITE
        for ci, val in enumerate([m, fn, pe]):
            _add_rect_text(sld, val, col_x[ci], y, col_w[ci], 0.44,
                           fill_color=bg, text_color=DARK,
                           font_size=10, line_color=GREY)


def slide_layer5(prs):
    layout = prs.slide_layouts[6]
    sld = prs.slides.add_slide(layout)
    _set_bg(sld, LIGHT)
    _section_header(sld, "Layer 5 — Dashboard & Reporting",
                    "Two dashboards: R Shiny (original) + Streamlit (Python live version)",
                    accent=_rgb((0x1e, 0x8b, 0x4c)))

    _add_rect_text(sld, "Streamlit Dashboard  (Python)", 0.3, 1.5, 4.4, 0.55,
                   fill_color=_rgb((0x1e,0x8b,0x4c)), font_size=13, bold=True)
    tabs = ["📊 Portfolio Overview", "⚖️ Weight Allocation", "📈 Analytics",
            "🛡️ Risk Monitor", "🔴 Live Data", "📰 Sentiment", "📡 Real-Time Signals"]
    for i, t in enumerate(tabs):
        _add_textbox(sld, t, 0.4, 2.1 + i * 0.44, 4.0, 0.42,
                     font_size=12, color=DARK)

    _add_rect_text(sld, "R Shiny + R Markdown  (original)", 5.1, 1.5, 4.6, 0.55,
                   fill_color=ORANGE, font_size=13, bold=True)
    r_items = ["Portfolio overview — equity curve", "Weight allocation — stacked area",
               "Sentiment heatmap (planned)", "Risk monitor — rolling VaR / CVaR",
               "R Markdown → HTML / PDF report", "Reads pre-computed backtest_results.rds"]
    for i, item in enumerate(r_items):
        _add_textbox(sld, "•  " + item, 5.2, 2.1 + i * 0.44, 4.4, 0.42,
                     font_size=12, color=DARK)

    _add_textbox(sld,
        "Both dashboards read the same pre-computed result files (CSVs + RDS).  "
        "Live Streamlit tab adds real-time yfinance refresh + LSTM re-inference.",
        0.3, 5.8, 9.4, 0.8, font_size=11, color=GREY)


def slide_why_r_python(prs):
    layout = prs.slide_layouts[6]
    sld = prs.slides.add_slide(layout)
    _set_bg(sld, LIGHT)
    _section_header(sld, "Why R + Python Together?",
                    "Each language does what it's genuinely best at")

    headers = ["Task", "Python", "R", "Winner"]
    col_x = [0.2, 2.8, 5.8, 8.5]; col_w = [2.5, 2.8, 2.5, 1.3]
    for ci, h in enumerate(headers):
        col = NAVY if h != "Winner" else GREEN
        _add_rect_text(sld, h, col_x[ci], 1.45, col_w[ci], 0.42,
                       fill_color=col, font_size=12, bold=True)

    rows = [
        ("Deep Learning / LSTM",         "✅ PyTorch",          "❌ keras (weaker)",      "Python"),
        ("Feature Engineering",          "✅ pandas, ta",       "⚠️  xts (workable)",     "Python"),
        ("Portfolio Constraint Solving",  "⚠️  cvxpy (manual)",  "✅ PortfolioAnalytics",  "R"),
        ("Risk Metrics (Sharpe/CVaR)",   "⚠️  manual NumPy",   "✅ PerformanceAnalytics","R"),
        ("Time-Series Backtesting",      "⚠️  manual loop",    "✅ xts + endpoints()",   "R"),
        ("Report Generation",            "⚠️  Jupyter/nbconv", "✅ R Markdown",           "R"),
        ("Interactive Dashboard",        "✅ Streamlit/Dash",   "✅ Shiny",               "Both"),
        ("Sentiment / NLP",              "✅ HuggingFace",      "❌ limited",             "Python"),
    ]
    py_col  = _rgb((0x2e, 0x86, 0xc1))
    r_col   = _rgb((0xe6, 0x7e, 0x22))
    win_col = _rgb((0x1e, 0x8b, 0x4c))

    for ri, (task, py, r, win) in enumerate(rows):
        y = 1.9 + ri * 0.52
        bg = LIGHT if ri % 2 == 0 else WHITE
        _add_rect_text(sld, task, col_x[0], y, col_w[0], 0.48,
                       fill_color=bg, text_color=DARK, font_size=11, line_color=GREY)
        _add_rect_text(sld, py,   col_x[1], y, col_w[1], 0.48,
                       fill_color=bg, text_color=py_col if "✅" in py else GREY,
                       font_size=10, line_color=GREY)
        _add_rect_text(sld, r,    col_x[2], y, col_w[2], 0.48,
                       fill_color=bg, text_color=r_col if "✅" in r else GREY,
                       font_size=10, line_color=GREY)
        wfill = win_col if win != "Both" else BLUE
        _add_rect_text(sld, win,  col_x[3], y, col_w[3], 0.48,
                       fill_color=wfill, text_color=WHITE, font_size=10)


def slide_results(prs):
    layout = prs.slide_layouts[6]
    sld = prs.slides.add_slide(layout)
    _set_bg(sld, NAVY)
    _add_rect(sld, 0, 0, SLD_W, 0.1, ORANGE)

    _add_textbox(sld, "Key Results — Test Period (Dec 2023 – Dec 2024)",
                 0.3, 0.15, 9.4, 0.75,
                 font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    cards = [
        ("DSR Constrained", "52.6%\nann. return", "Sharpe  2.51", "Max DD  14.3%", BLUE),
        ("DSR Unconstrained", "69.1%\nann. return", "Sharpe  3.02", "Max DD  12.9%", _rgb((0x8e,0x44,0xad))),
        ("Equal Weight\n(benchmark)", "43.2%\nann. return", "Sharpe  2.13", "Max DD  14.4%", GREY),
    ]
    for i, (title, ret, sharpe, mdd, col) in enumerate(cards):
        x = 0.4 + i * 3.2
        _add_rect_text(sld, title, x, 1.1, 2.8, 0.6,
                       fill_color=col, font_size=13, bold=True)
        _add_rect_text(sld, ret, x, 1.75, 2.8, 1.1,
                       fill_color=WHITE, text_color=col, font_size=22, bold=True)
        _add_rect_text(sld, sharpe, x, 2.9, 2.8, 0.6,
                       fill_color=LIGHT, text_color=DARK, font_size=14)
        _add_rect_text(sld, mdd, x, 3.55, 2.8, 0.6,
                       fill_color=LIGHT, text_color=RED, font_size=14)

    _add_textbox(sld,
        "DSR Constrained beats Equal Weight by +9.4 pp annualised return  "
        "with a meaningfully higher Sharpe ratio and similar drawdown.",
        0.4, 4.35, 9.2, 0.75, font_size=14, color=LBLUE, align=PP_ALIGN.CENTER)

    _add_textbox(sld,
        "Model trained: 2018–2024  ·  5 assets  ·  60-day lookback  ·  "
        "LSTM (128→64)  ·  DSR loss (η=0.01)  ·  PortfolioAnalytics constraints",
        0.4, 5.15, 9.2, 0.65, font_size=11, color=GREY, align=PP_ALIGN.CENTER)

    _add_rect(sld, 0, SLD_H - 0.1, SLD_W, 0.1, ORANGE)


def slide_bear_market(prs):
    layout = prs.slide_layouts[6]
    sld = prs.slides.add_slide(layout)
    _set_bg(sld, NAVY)
    _add_rect(sld, 0, 0, SLD_W, 0.1, ORANGE)

    _add_textbox(sld, "Stress Test: Does it hold in a Bear Market?",
                 0.3, 0.15, 9.4, 0.7,
                 font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    _add_textbox(sld, "Retrained on 2018–2021 only  ·  Tested on 2022 (out-of-sample crash year)",
                 0.3, 0.82, 9.4, 0.4,
                 font_size=12, color=GREY, align=PP_ALIGN.CENTER)

    # ── Matplotlib bar chart ─────────────────────────────────────────────────
    fig, axes = plt.subplots(1, 2, figsize=(9.5, 3.8))
    fig.patch.set_facecolor(mNAVY)

    datasets = [
        ("Bull Market 2024\n(Dec 2023 – Dec 2024)",
         ["DSR Constrained", "Equal Weight"],
         [52.6, 43.2],
         mBLUE, "#5d6d7e"),
        ("Bear Market 2022\n(Jan – Dec 2022)",
         ["DSR Constrained", "Equal Weight"],
         [-34.3, -39.8],
         mBLUE, "#5d6d7e"),
    ]

    for ax, (title, labels, vals, c1, c2) in zip(axes, datasets):
        colors = [c1, c2]
        bars = ax.bar(labels, vals, color=colors, width=0.5, zorder=3)
        ax.set_facecolor(mNAVY)
        ax.tick_params(colors="white", labelsize=9)
        ax.spines["bottom"].set_color("#aed6f1")
        ax.spines["left"].set_color("#aed6f1")
        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)
        ax.yaxis.label.set_color("white")
        ax.set_ylabel("Ann. Return (%)", color="white", fontsize=9)
        ax.set_title(title, fontsize=10, fontweight="bold", color="white", pad=6)
        ax.axhline(0, color="#aed6f1", lw=0.8, linestyle="--")

        for bar, val in zip(bars, vals):
            label = f"{val:+.1f}%"
            ypos = val + (1.5 if val >= 0 else -3.5)
            ax.text(bar.get_x() + bar.get_width()/2, ypos,
                    label, ha="center", va="bottom", fontsize=11,
                    fontweight="bold", color="white")

        # Outperformance annotation
        diff = vals[0] - vals[1]
        ax.annotate(f"↑ +{diff:.1f}pp\nvs benchmark",
                    xy=(0.5, 0.92), xycoords="axes fraction",
                    ha="center", fontsize=9, color=mORANGE, fontweight="bold")

    plt.tight_layout(pad=1.2)
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                facecolor=mNAVY)
    buf.seek(0)
    plt.close(fig)
    sld.shapes.add_picture(buf, Inches(0.25), Inches(1.35), Inches(9.5), Inches(3.9))

    # ── Takeaway banner ──────────────────────────────────────────────────────
    _add_rect(sld, 0.3, 5.4, 9.4, 0.6, GREEN)
    _add_textbox(sld,
        "Model consistently outperforms the equal-weight benchmark "
        "in both bull (+9.4 pp) and bear (+5.5 pp) markets",
        0.3, 5.4, 9.4, 0.6,
        font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    _add_rect(sld, 0, SLD_H - 0.1, SLD_W, 0.1, ORANGE)

    # ── Thank you footer ─────────────────────────────────────────────────────
    _add_textbox(sld, "Thank you",
                 0.3, 6.1, 9.4, 0.5,
                 font_size=16, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width  = Inches(SLD_W)
    prs.slide_height = Inches(SLD_H)

    print("Building slides...")
    slide_title(prs);                      print("  1/16  Title")
    slide_what_are_we_building(prs);       print("  2/16  What are we building?")
    slide_architecture(prs);              print("  3/16  System architecture diagram")
    slide_data_flow(prs);                 print("  4/16  Data flow diagram")
    slide_layer1(prs);                    print("  5/16  Layer 1 — Data ingestion")
    slide_layer2(prs);                    print("  6/16  Layer 2 — Feature engineering")
    slide_layer3(prs);                    print("  7/16  Layer 3 — LSTM + DSR")
    slide_r_overview(prs);                print("  8/16  R overview — component diagram")
    slide_reticulate(prs);                print("  9/16  R-Python sequence diagram")
    slide_backtest(prs);                  print(" 10/16  Walk-forward backtest")
    slide_cvar(prs);                      print(" 11/16  CVaR constraint flow")
    slide_performance_analytics(prs);     print(" 12/16  PerformanceAnalytics metrics")
    slide_layer5(prs);                    print(" 13/16  Dashboard & reporting")
    slide_why_r_python(prs);              print(" 14/16  Why R + Python?")
    slide_results(prs);                   print(" 15/16  Bull market results")
    slide_bear_market(prs);               print(" 16/16  Bear market stress test")

    out = "presentation.pptx"
    prs.save(out)
    print(f"\nSaved → {out}")


if __name__ == "__main__":
    main()
